""""""
from genericpath import exists
import pptx
from os.path import exists

import logging

logger = logging.getLogger(__name__)

class PowerPointParser:
    """
    
    ---
    Attributes
    ----------


    Methods
    -------

    
    """

    name = "PowerPoint Parser"
    version = "0.0.1"

    def __init__(self):
        pass

    def __str__(self):
        return "{} {}".format(self.name, self.version)
        
    def get_extensions(self):
        """
        
        Parameters
        ----------

        Returns
        -------
        None
        """
        return [".pptx"]

    def get_functions(self):
        """
        
        Parameters
        ----------

        Returns
        -------
        None
        """
        return {
            "metadata": self.get_metadata,
            "analyse": self.analyze,
            "contents": self.get_contents
        }

    def get_metadata(self, filepath: str) -> dict:
        """This will return the doc info infomation from the 
        Named file.
        
        Parameters
        ----------

        Returns
        -------
        None
        """

        data = {}
    
        if not filepath or not exists(filepath):
            return data

        doc = pptx.Presentation(filepath)
        
        # get the core properties from the file...
        # https://python-docx.readthedocs.io/en/latest/api/document.html#coreproperties-objects
        cp = doc.core_properties
        
        data['author'] = cp.author
        data['category'] = cp.category
        data['comments'] = cp.comments
        data['content_status'] = cp.content_status
        data['created'] = cp.created
        data['identifier'] = cp.identifier
        data['keywords'] = cp.keywords
        data['language'] = cp.language
        data['last_modified_by'] = cp.last_modified_by
        data['last_printed'] = cp.last_printed
        data['modified'] = cp.modified
        data['revision'] = cp.revision
        data['subject'] = cp.subject
        data['title'] = cp.title
        data['version'] = cp.version
                
        return data

    def analyze(self, filepath: str) -> dict:
        """
        
        Parameters
        ----------

        Returns
        -------
        None
        """
        data = dict

        if not filepath or not exists(filepath):
            return data

        
        return data

    def get_contents(self, filepath: str) -> list:
        """This will return the paragraph objects in a word document
        
        Parameters
        ----------

        Returns
        -------
        None
        """
        data = []

        if not filepath or not exists(filepath):
            return data

        prs = pptx.Presentation(filepath)
        for idx, slide in enumerate(prs.slides):
            data.append("Page: {}".format(idx))
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if len(shape.text) > 0:
                        data.append(shape.text)
                elif shape.has_table:
                    data.append("== table")
                    for r in shape.table.rows:
                        data.append("== row")
                        for c in r.cells:
                            data.append(c.text)
                    data.append("== end")
        return {
            "type": "strings",
            "strings": data
        }
        

if __name__ == "__main__":
    from argparse import ArgumentParser
    argparser = ArgumentParser(
        prog="Visio Scanner",
        description="will scan a visio file and return the discovered information")

    argparser.add_argument('-fp', '--file_path',
        dest="file_path",
        help="The path for the file to scan.") 

    args = argparser.parse_args()

    fp = args.file_path
    if not fp:
        #fp = r"E:\users\ms\google\thoughtswin\Manitoba Hydro\General\Reference\Strategy 2040 Leadership Town Hall Deck.pptx"
        fp = r"E:\users\ms\google\thoughtswin\Manitoba Hydro\General\Reference\EA Info.pptx"

    obj = PowerPointParser()
    
    data = obj.get_metadata(filepath=fp)
    logger.info(data)

    data = obj.get_contents(filepath=fp)
    logger.info(data)

    data = obj.analyze(filepath=fp)
    logger.info(data)        